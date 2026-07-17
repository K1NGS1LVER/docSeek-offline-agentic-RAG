"""
End-to-end tests against a real docSeek server over HTTP.

Covers the full agentic RAG surface: ingestion with every chunking strategy,
hybrid search, cross-encoder reranking, the agentic /ask SSE protocol
(trace/sources/answer events), document view, deletion, and rebuild.

LLM-dependent behavior degrades to heuristics when Ollama is down, so these
tests assert pipeline structure (stages, events, ordering), not which planner
produced it — the suite passes with or without Ollama running.
"""

import json

import pytest
import requests

# Read timeout is generous: the first rerank lazily loads the cross-encoder,
# and agent planning/grading calls a local LLM.
TIMEOUT = (10, 180)

DOC_AGENT = """docSeek is a local-first retrieval augmented generation system. All data stays on the user's device. The system never sends documents to external services. Privacy is the core design principle of the whole application.

The retrieval agent plans each query before searching. It classifies the question and chooses how many chunks to retrieve. It can rewrite unclear queries into better search terms. For complex questions it splits the query into simpler subqueries.

The grading step judges whether retrieved evidence answers the question. If the evidence is weak the agent reformulates the query and retries. The loop is bounded to keep latency predictable."""

DOC_PLANETS = (
    "The solar system contains eight planets orbiting the sun. "
    "Jupiter is the largest planet with a mass greater than all others combined. "
    "Saturn is famous for its ring system made of ice and rock. "
    "Mars has the tallest volcano known as Olympus Mons. "
    "Venus has a runaway greenhouse atmosphere hotter than Mercury. "
) * 2 + "\n\n" + (
    "Sourdough bread requires a fermented starter culture of wild yeast. "
    "The dough must be folded several times during bulk fermentation. "
    "A dutch oven traps steam and produces a crisp crust. "
    "Long cold retardation in the fridge develops complex flavor. "
    "Scoring the loaf controls how the bread expands in the oven. "
) * 2


def upload(server, filename, content, strategy=None):
    files = {"file": (filename, content.encode())}
    data = {"chunking_strategy": strategy} if strategy else {}
    return requests.post(f"{server}/upload", files=files, data=data, timeout=TIMEOUT)


def upload_bytes(server, filename, data, content_type):
    """Upload raw binary content (for PDF/PPTX/DOCX)."""
    files = {"file": (filename, data, content_type)}
    return requests.post(f"{server}/upload", files=files, timeout=TIMEOUT)


def make_pdf_bytes(text):
    """A minimal but spec-correct single-page text PDF (no writer dependency)."""
    objs = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[3 0 R]/Count 1>>",
        b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R"
        b"/Resources<</Font<</F1 5 0 R>>>>>>",
    ]
    stream = b"BT /F1 24 Tf 72 700 Td (" + text.encode() + b") Tj ET"
    objs.append(b"<</Length " + str(len(stream)).encode() + b">>stream\n" + stream + b"\nendstream")
    objs.append(b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>")
    out = b"%PDF-1.4\n"
    offsets = []
    for i, body in enumerate(objs, 1):
        offsets.append(len(out))
        out += str(i).encode() + b" 0 obj" + body + b"\nendobj\n"
    xref_pos = len(out)
    n = len(objs) + 1
    out += b"xref\n0 " + str(n).encode() + b"\n0000000000 65535 f \n"
    for off in offsets:
        out += ("%010d 00000 n \n" % off).encode()
    out += (b"trailer<</Size " + str(n).encode() + b"/Root 1 0 R>>\nstartxref\n"
            + str(xref_pos).encode() + b"\n%%EOF")
    return out


def make_pptx_bytes(lines):
    """A one-slide PowerPoint deck with the given text lines."""
    import io
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4)).text_frame
    tf.text = lines[0]
    for ln in lines[1:]:
        tf.add_paragraph().text = ln
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def read_sse(resp, stop_after_event=None, stop_after_answer_chars=None):
    """Parse an SSE stream into (events, answer_text).

    events is a list of (event_name, parsed_payload); unnamed events are
    "message" and their string payloads accumulate into answer_text.
    Closing early (via the stop_* args) is deliberate: it verifies streaming
    without waiting for a full local-LLM generation.
    """
    events = []
    answer = ""
    name = ""
    data_lines = []
    try:
        for raw in resp.iter_lines(decode_unicode=True):
            if raw:
                if raw.startswith("event:"):
                    name = raw[6:].strip()
                elif raw.startswith("data:"):
                    data_lines.append(raw[5:].lstrip(" "))
                continue
            # Blank line terminates one SSE event block.
            if not data_lines:
                name = ""
                continue
            payload = "\n".join(data_lines)
            try:
                parsed = json.loads(payload)
            except json.JSONDecodeError:
                parsed = payload
            event_name = name or "message"
            events.append((event_name, parsed))
            if event_name == "message" and isinstance(parsed, str):
                answer += parsed
            name, data_lines = "", []
            if stop_after_event and event_name == stop_after_event:
                break
            if stop_after_answer_chars and len(answer) >= stop_after_answer_chars:
                break
    finally:
        resp.close()
    return events, answer


@pytest.fixture(scope="session")
def corpus(server):
    """Ingest the shared test corpus once; returns upload responses."""
    agent_doc = upload(server, "agent_doc.txt", DOC_AGENT)
    planets_doc = upload(server, "planets_and_bread.txt", DOC_PLANETS, strategy="semantic")
    assert agent_doc.status_code == 200, agent_doc.text
    assert planets_doc.status_code == 200, planets_doc.text
    return {"agent_doc": agent_doc.json(), "planets_doc": planets_doc.json()}


# ---------------------------------------------------------------- system


def test_stats_reports_agentic_config(server):
    stats = requests.get(f"{server}/stats", timeout=TIMEOUT).json()
    assert stats["agentic"] is True
    assert stats["hybrid_search"] is True
    assert "cross-encoder" in stats["reranker_model"]
    assert stats["chunking_strategy"] == "auto"


def test_ask_with_empty_corpus(server):
    """Must run before any ingestion: the corpus fixture is instantiated by
    the first test that requests it, and pytest executes in file order."""
    resp = requests.post(
        f"{server}/ask", json={"query": "anything"}, stream=True, timeout=TIMEOUT
    )
    _, answer = read_sse(resp)
    assert "No documents have been uploaded yet" in answer


# ------------------------------------------------------------- ingestion


def test_upload_auto_strategy(corpus):
    doc = corpus["agent_doc"]
    assert doc["status"] == "success"
    assert doc["chunks"] >= 1
    assert doc["chunking"] in ("recursive", "semantic")
    assert len(doc["doc_ids"]) == doc["chunks"]


def test_upload_semantic_strategy(corpus):
    doc = corpus["planets_doc"]
    assert doc["chunking"] == "semantic"
    # Two clearly distinct topics must not end up in a single chunk.
    assert doc["chunks"] >= 2


def test_upload_rejects_unknown_strategy(server):
    r = upload(server, "bad_strategy.txt", "some text", strategy="quantum")
    assert r.status_code == 400
    assert "chunking strategy" in r.json()["detail"].lower()


def test_upload_rejects_unsupported_extension(server):
    r = upload(server, "binary.exe", "not really a doc")
    assert r.status_code == 400


def test_upload_multiple(server):
    files = [
        ("files", ("multi_a.txt", b"Alpha document about testing multiple uploads.")),
        ("files", ("multi_b.txt", b"Beta document about testing multiple uploads.")),
    ]
    r = requests.post(f"{server}/upload-multiple", files=files, timeout=TIMEOUT)
    assert r.status_code == 200
    body = r.json()
    assert body["files_processed"] == 2
    assert all(item["status"] == "success" for item in body["results"])


def test_concurrent_uploads_keep_index_consistent(server):
    """Parallel uploads must not corrupt the FAISS index: the ingest lock keeps
    DB inserts and index adds in lockstep, so total_documents stays equal to
    total_vectors even under concurrency."""
    from concurrent.futures import ThreadPoolExecutor

    def one(n):
        return upload(
            server, f"concurrent_{n}.txt",
            f"Concurrent upload number {n} about widget calibration and gear ratios.",
        )

    with ThreadPoolExecutor(max_workers=6) as pool:
        responses = list(pool.map(one, range(6)))

    assert all(r.status_code == 200 for r in responses), [r.text for r in responses]

    stats = requests.get(f"{server}/stats", timeout=TIMEOUT).json()
    assert stats["total_documents"] == stats["total_vectors"], stats


def test_ingest_raw_text(server):
    r = requests.post(
        f"{server}/ingest",
        json={"text": "Raw ingested snippet about xylophone maintenance."},
        timeout=TIMEOUT,
    )
    assert r.status_code == 200
    assert r.json()["id"] > 0


def test_upload_pdf_extracts_and_indexes(server):
    data = make_pdf_bytes("Photosynthesis converts sunlight into chemical energy in plants.")
    r = upload_bytes(server, "biology.pdf", data, "application/pdf")
    assert r.status_code == 200, r.text
    body = r.json()
    assert body["status"] == "success" and body["chunks"] >= 1

    hits = requests.post(
        f"{server}/search", json={"query": "photosynthesis sunlight energy", "k": 3}, timeout=TIMEOUT
    ).json()
    assert any("Photosynthesis" in h["content"] for h in hits)


def test_upload_pptx_extracts_and_indexes(server):
    data = make_pptx_bytes([
        "Quarterly Business Review",
        "Revenue grew twenty percent this quarter.",
        "Cloud adoption accelerated across every region.",
    ])
    r = upload_bytes(
        server, "deck.pptx", data,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    assert r.status_code == 200, r.text
    body = r.json()
    assert body["status"] == "success" and body["chunks"] >= 1

    hits = requests.post(
        f"{server}/search", json={"query": "quarterly revenue growth cloud", "k": 3}, timeout=TIMEOUT
    ).json()
    assert any("Revenue grew" in h["content"] for h in hits)


def test_upload_scanned_pdf_ocr_or_clean_failure(server):
    """A scanned (image-only) PDF has no text layer. With the OCR stack present
    it ingests via the Tesseract fallback; without it, ingestion fails cleanly
    (400, no content) rather than 500. Structural, like the other optional-model
    tests -- passes with or without OCR installed."""
    import io
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (1400, 220), "white")
    draw = ImageDraw.Draw(img)
    font = None
    for path in (
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ):
        try:
            font = ImageFont.truetype(path, 60)
            break
        except Exception:
            pass
    draw.text((30, 70), "Photosynthesis powers plant growth",
              fill="black", font=font or ImageFont.load_default())
    buf = io.BytesIO()
    img.save(buf, "PDF", resolution=150)

    r = upload_bytes(server, "scanned.pdf", buf.getvalue(), "application/pdf")
    assert r.status_code in (200, 400), r.text
    if r.status_code == 200:
        assert r.json()["chunks"] >= 1


# --------------------------------------------------------------- search


def test_search_hybrid_finds_content(server, corpus):
    r = requests.post(
        f"{server}/search",
        json={"query": "which planet is the largest", "k": 3},
        timeout=TIMEOUT,
    )
    assert r.status_code == 200
    results = r.json()
    assert results, "expected at least one hit"
    assert any("Jupiter" in res["content"] for res in results)
    for res in results:
        assert set(res) >= {"id", "score", "content", "source"}


def test_search_rerank_scores_and_orders(server, corpus):
    r = requests.post(
        f"{server}/search",
        json={"query": "how does the agent decide what to retrieve", "k": 3, "rerank": True},
        timeout=TIMEOUT,
    )
    assert r.status_code == 200
    results = r.json()
    assert results
    scores = [res["rerank_score"] for res in results]
    assert all(s is not None for s in scores)
    assert scores == sorted(scores, reverse=True)


def test_search_scores_respect_threshold_contract(server, corpus):
    """Weak dense-only hits are dropped; keyword-only hits surface score 0.0.
    An unrelated query may or may not clear the threshold (embedding
    similarity of arbitrary text is not bounded), so assert the score
    contract rather than emptiness."""
    r = requests.post(
        f"{server}/search",
        json={"query": "zeppelin cactus harmonica blizzard", "k": 5},
        timeout=TIMEOUT,
    )
    assert r.status_code == 200
    for res in r.json():
        assert res["score"] == 0.0 or res["score"] >= 0.20


# ------------------------------------------------------------------ ask


def test_ask_agentic_streams_trace_and_sources(server, corpus):
    resp = requests.post(
        f"{server}/ask",
        json={"query": "How does docSeek protect user privacy?"},
        stream=True,
        timeout=TIMEOUT,
    )
    assert resp.status_code == 200
    events, _ = read_sse(resp, stop_after_event="sources")

    stages = [ev[1]["stage"] for ev in events if ev[0] == "trace"]
    assert "plan" in stages
    assert "retrieve" in stages
    assert "grade" in stages
    # Stages must appear in pipeline order.
    assert stages.index("plan") < stages.index("retrieve") < stages.index("grade")

    sources = [ev[1] for ev in events if ev[0] == "sources"]
    assert len(sources) == 1
    assert 1 <= len(sources[0]) <= 12  # within AGENT_MAX_K
    assert any("privacy" in s["content"].lower() for s in sources[0])
    for s in sources[0]:
        assert set(s) >= {"id", "content", "score", "source"}


def test_ask_agentic_respects_explicit_k(server, corpus):
    resp = requests.post(
        f"{server}/ask",
        json={"query": "what is jupiter", "k": 2},
        stream=True,
        timeout=TIMEOUT,
    )
    events, _ = read_sse(resp, stop_after_event="sources")
    sources = [ev[1] for ev in events if ev[0] == "sources"]
    assert len(sources) == 1
    assert len(sources[0]) <= 2

    plans = [ev[1] for ev in events if ev[0] == "trace" and ev[1]["stage"] == "plan"]
    assert plans and plans[0]["data"]["k"] == 2


def test_ask_non_agentic_has_sources_but_no_trace(server, corpus, ollama_up):
    resp = requests.post(
        f"{server}/ask",
        json={"query": "what makes sourdough bread rise", "k": 2, "agentic": False},
        stream=True,
        timeout=TIMEOUT,
    )
    # Read a few answer characters to prove deltas stream, then hang up.
    events, answer = read_sse(resp, stop_after_answer_chars=10)
    assert not [ev for ev in events if ev[0] == "trace"]
    sources = [ev[1] for ev in events if ev[0] == "sources"]
    assert len(sources) == 1 and len(sources[0]) >= 1
    assert answer, "expected streamed answer text (or an inline Ollama error message)"
    if ollama_up:
        assert "Error communicating with Ollama" not in answer


def test_ask_always_emits_sources_event(server, corpus):
    """The sources event is part of the wire contract for every non-empty
    corpus, whether or not retrieval found anything for this query."""
    resp = requests.post(
        f"{server}/ask",
        json={"query": "quokka submarine oboe glacier", "agentic": False, "k": 3},
        stream=True,
        timeout=TIMEOUT,
    )
    events, answer = read_sse(resp, stop_after_answer_chars=10)
    sources = [ev[1] for ev in events if ev[0] == "sources"]
    assert len(sources) == 1
    if sources[0] == []:
        assert "couldn't find any relevant information" in answer
    else:
        assert answer


# -------------------------------------------------- document view / delete


def test_document_view_highlights_chunk(server, corpus):
    chunk_id = corpus["agent_doc"]["doc_ids"][0]
    r = requests.get(f"{server}/document/view", params={"id": chunk_id}, timeout=TIMEOUT)
    assert r.status_code == 200
    assert 'id="target"' in r.text
    assert "agent_doc.txt" in r.text


def test_delete_documents_by_source(server):
    up = upload(server, "to_delete.txt", "Ephemeral document about disposable content.")
    assert up.status_code == 200
    doc_ids = up.json()["doc_ids"]

    view = requests.get(f"{server}/document/view", params={"id": doc_ids[0]}, timeout=TIMEOUT)
    assert view.status_code == 200

    # Resolve the stored source_file path via search metadata.
    hit = requests.post(
        f"{server}/search",
        json={"query": "ephemeral disposable content", "k": 1},
        timeout=TIMEOUT,
    ).json()
    assert hit
    source_file = hit[0]["source"]["source_file"]

    r = requests.delete(
        f"{server}/documents", params={"source_file": source_file}, timeout=TIMEOUT
    )
    assert r.status_code == 200
    assert r.json()["db_rows"] == len(doc_ids)

    gone = requests.get(f"{server}/document/view", params={"id": doc_ids[0]}, timeout=TIMEOUT)
    assert gone.status_code == 404


# --------------------------------------------------------------- rebuild


def test_rebuild_reindexes_everything(server, corpus):
    stats_before = requests.get(f"{server}/stats", timeout=TIMEOUT).json()
    r = requests.post(f"{server}/rebuild", timeout=TIMEOUT)
    assert r.status_code == 200
    assert r.json()["documents_indexed"] == stats_before["total_documents"]

    # Search must still work against the rebuilt index.
    hits = requests.post(
        f"{server}/search",
        json={"query": "largest planet", "k": 3},
        timeout=TIMEOUT,
    ).json()
    assert any("Jupiter" in h["content"] for h in hits)
