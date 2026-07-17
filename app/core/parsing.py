import io
import logging
import re
from typing import List, Tuple
from bs4 import BeautifulSoup

logger = logging.getLogger(__name__)

CHUNK_SIZE = 1000  # Characters per chunk (~200 tokens, well under mpnet's 384)
CHUNK_OVERLAP = 150  # Overlap between chunks

def chunk_text(
    text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP
) -> List[Tuple[str, int, int]]:
    """
    Split text into overlapping chunks with character positions.
    Returns list of (chunk_text, start_char, end_char) tuples.
    This prevents context loss at chunk boundaries.
    """
    if not text:
        return []
        
    if len(text) <= chunk_size:
        return [(text, 0, len(text))]

    chunks = []
    start = 0

    while start < len(text):
        end = min(start + chunk_size, len(text))
        original_end = end

        # Try to break at sentence boundary
        if end < len(text):
            search_start = max(0, chunk_size // 2) 
            text_slice = text[start:end]
            
            for punct in [". ", "! ", "? ", "\n\n"]:
                last_punct = text_slice.rfind(punct, search_start) 
                if last_punct != -1:
                    end = start + last_punct + len(punct)
                    break
            
        chunk = text[start:end].strip()
        if chunk:
            # Find actual positions after stripping
            chunk_start = text.find(chunk, start)
            chunk_end = chunk_start + len(chunk)
            chunks.append((chunk, chunk_start, chunk_end))

        next_start = end - overlap
        if next_start <= start:
            next_start = start + max(1, chunk_size - overlap)
            
        start = next_start

    return chunks

def parse_markdown(content: str) -> str:
    """Clean markdown content"""
    # Remove front matter (YAML)
    content = re.sub(r"^---\n.*?\n---\n", "", content, flags=re.DOTALL)
    return content

def parse_html(content: str) -> str:
    """Parse HTML content and extract text"""
    soup = BeautifulSoup(content, "html.parser")
    for script in soup(["script", "style", "nav", "footer"]):
        script.decompose()
    return soup.get_text(separator="\n", strip=True)

def parse_pdf(data: bytes) -> str:
    """Extract text from a PDF, page by page.

    Uses pypdf (pure Python). Text-layer PDFs extract cleanly; scanned/image
    PDFs have no text layer and yield little or nothing (no OCR here).
    """
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            pages.append(text.strip())
    combined = clean_text("\n\n".join(pages))

    # No text layer (scanned/image-only PDF): fall back to local OCR if available.
    if not combined.strip():
        from . import ocr

        ocr_text = ocr.ocr_pdf(data)
        if ocr_text.strip():
            logger.info("PDF had no text layer; extracted text via local OCR fallback.")
            return clean_text(ocr_text)
    return combined


def parse_pptx(data: bytes) -> str:
    """Extract text from a PowerPoint deck: per-slide shape text and tables."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_lines.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_lines.append(" | ".join(cells))
        if slide_lines:
            parts.append("\n".join(slide_lines))
    # Blank line between slides keeps them as separate semantic blocks.
    return clean_text("\n\n".join(parts))


def clean_text(content: str) -> str:
    """General text cleaning"""
    content = re.sub(r"\n{3,}", "\n\n", content)
    return content.strip()
